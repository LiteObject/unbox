"""PowerPoint (.pptx) text extractor using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from unbox.base import BaseExtractor


class PptxExtractor(BaseExtractor):
    """Extract plain text from PowerPoint presentations."""

    supported_extensions = [".pptx"]

    def extract(self, file_path: Path) -> str:
        """Extract text from all slides of a PowerPoint presentation.

        Parameters
        ----------
        file_path:
            Path to the ``.pptx`` file.

        Returns
        -------
        str
            Slide text separated by blank lines, with a heading per slide.
        """
        prs = Presentation(str(file_path))
        slides_text: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            if len(parts) > 1:  # more than just the header
                slides_text.append("\n".join(parts))

        return "\n\n".join(slides_text)
